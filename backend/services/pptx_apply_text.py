from __future__ import annotations

import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.text.text import TextFrame

from backend.services.font_manager import clone_font_props, contains_cjk
from backend.services.pptx_xml_core import get_pptx_theme_summary


# Pre-compile Regex for CJK Kinsoku optimization
# Match CJK char, space, CJK char -> remove space
CJK_SPACE_PATTERN = re.compile(
    r'([\u4e00-\u9fff\u3400-\u4dbf\u3000-\u303f])\s+([\u4e00-\u9fff\u3400-\u4dbf\u3000-\u303f])'
)

# Remove characters that are invalid in XML 1.0 to avoid python-pptx text failures.
INVALID_XML_CHARS = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")


def sanitize_xml_text(text: str) -> str:
    if not text:
        return text
    return INVALID_XML_CHARS.sub("", text)


def capture_full_frame_styles(text_frame: TextFrame) -> list[dict[str, Any]]:
    """Captures detailed style info for every paragraph and its first run."""
    styles = []
    for p in text_frame.paragraphs:
        p_style = {
            "level": p.level,
            "alignment": p.alignment,
            "space_before": p.space_before,
            "space_after": p.space_after,
            "line_spacing": p.line_spacing,
            "font_obj": p.runs[0].font if p.runs else None
        }
        styles.append(p_style)
    return styles


def apply_paragraph_style(
    paragraph, 
    p_style: dict[str, Any], 
    scale: float = 1.0, 
    target_language: str | None = None,
    font_mapping: dict[str, list[str]] | None = None
) -> None:
    """Applies paragraph-level settings and clones font from font_obj."""
    try:
        paragraph.level = p_style.get("level", 0)
        paragraph.alignment = p_style.get("alignment")
        if p_style.get("space_before") is not None:
            paragraph.space_before = p_style["space_before"]
        if p_style.get("space_after") is not None:
            paragraph.space_after = p_style["space_after"]
        if p_style.get("line_spacing") is not None:
            paragraph.line_spacing = p_style["line_spacing"]
            
        source_font = p_style.get("font_obj")
        if source_font and paragraph.runs:
            clone_font_props(source_font, paragraph.runs[0].font, target_language=target_language, font_mapping=font_mapping)
            if scale != 1.0 and paragraph.runs[0].font.size:
                paragraph.runs[0].font.size = int(paragraph.runs[0].font.size * scale)
    except Exception:
        pass


def apply_shape_highlight(
    shape: Any,
    fill_color: RGBColor | None = None,
    line_color: RGBColor | None = None,
    dash_style: MSO_LINE_DASH_STYLE | None = None,
) -> None:
    """Apply highlighting styles to a shape.
    
    Args:
        shape: The shape object to modify.
        fill_color: Optional fill color (RGBColor).
        line_color: Optional line color (RGBColor).
        dash_style: Optional line dash style (MSO_LINE_DASH_STYLE).
    """
    if fill_color:
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        except Exception:
            pass  # Some shapes might not support fill

    if line_color:
        try:
            shape.line.color.rgb = line_color
        except Exception:
            pass

    if dash_style:
        try:
            shape.line.dash_style = dash_style
        except Exception:
            pass


def set_text_preserve_format(text_frame: TextFrame, new_text: str, auto_size: bool = False, scale: float = 1.0) -> None:
    try:
        new_text = sanitize_xml_text(new_text)
        para_styles = capture_full_frame_styles(text_frame)
        
        if auto_size:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
        text_frame.clear()
        lines = new_text.split("\n")
        
        for index, line in enumerate(lines):
            if index == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = line
            style_idx = min(index, len(para_styles) - 1) if para_styles else -1
            if style_idx >= 0:
                apply_paragraph_style(paragraph, para_styles[style_idx], scale=scale)
                
    except Exception:
        try:
            if auto_size:
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            text_frame.clear()
            text_frame.text = new_text
        except Exception:
            pass


def parse_hex_color(value: str | None, default: RGBColor) -> RGBColor:
    if not value:
        return default
    cleaned = value.strip().lstrip("#")
    if len(cleaned) != 6:
        return default
    try:
        return RGBColor.from_string(cleaned.upper())
    except ValueError:
        return default
    return default


def parse_dash_style(value: str | None) -> MSO_LINE_DASH_STYLE | None:
    if not value:
        return None
    normalized = value.strip().lower()
    mapping = {
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dashdot": MSO_LINE_DASH_STYLE.DASH_DOT,
        "solid": MSO_LINE_DASH_STYLE.SOLID,
    }
    return mapping.get(normalized)


def apply_cjk_line_breaking(text: str) -> str:
    """Simple implementation of Kinsoku Shori using compiled regex."""
    if not contains_cjk(text):
        return text
    
    # Use pre-compiled regex to remove spaces between CJK characters
    return CJK_SPACE_PATTERN.sub(r'\1\2', text)


def build_corrected_lines(source_text: str, translated_text: str) -> list[str]:
    source_lines = source_text.split("\n")
    non_cjk_lines = []
    for line in source_lines:
        if line.strip() == "":
            non_cjk_lines.append("")
            continue
        if contains_cjk(line):
            continue
        non_cjk_lines.append(line)
    while non_cjk_lines and non_cjk_lines[-1] == "":
        non_cjk_lines.pop()
    translated_lines = translated_text.split("\n") if translated_text else []  
    
    # optimize usage of apply_cjk_line_breaking
    optimized_lines = []
    for line in translated_lines:
        optimized_lines.append(apply_cjk_line_breaking(line))
        
    if non_cjk_lines:
        return non_cjk_lines + [""] + optimized_lines
    return optimized_lines


def set_bilingual_text(
    text_frame: TextFrame,
    source_text: str,
    translated_text: str,
    auto_size: bool = False,
    scale: float = 1.0,
    theme_data: dict[str, Any] | None = None,
    target_language: str | None = None,
    font_mapping: dict[str, list[str]] | None = None,
) -> bool:
    # Default to a nice blue if theme color accent1 is not available
    translated_color = RGBColor(0x1F, 0x77, 0xB4)
    source_text = sanitize_xml_text(source_text)
    translated_text = sanitize_xml_text(translated_text)
    if theme_data and 'colors' in theme_data:
        accent1 = theme_data['colors'].get('accent1')
        if accent1:
            try:
                translated_color = RGBColor.from_string(accent1)
            except Exception:
                pass
    
    text_scale = scale  # Use the calculated scale

    try:
        # 1. Capture original paragraph styles (including font objects)
        para_styles = capture_full_frame_styles(text_frame)

        if auto_size:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        text_frame.clear()
        
        # Add Source Text
        source_lines = source_text.split("\n")
        for index, line in enumerate(source_lines):
            if index == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = line
            style_idx = min(index, len(para_styles) - 1) if para_styles else -1
            if style_idx >= 0:
                # 2026/01/20: Apply scaling to source text too, to prevent bilingual overflow
                apply_paragraph_style(paragraph, para_styles[style_idx], scale=text_scale)

        # Add Empty Line for Separation
        separator = text_frame.add_paragraph()
        separator.text = " "
        base_size = 120000
        if para_styles and para_styles[0].get("font_obj"):
            base_size = para_styles[0]["font_obj"].size or base_size
        separator.font.size = int(base_size * 0.5)

        # Add Translated Text
        translated_lines = translated_text.split("\n")
        
        # Apply CJK optimization
        translated_lines = [apply_cjk_line_breaking(line) for line in translated_lines]
        
        for index, line in enumerate(translated_lines):
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            
            style_idx = min(index, len(para_styles) - 1) if para_styles else -1
            if style_idx >= 0:
                # Apply base style and then override color
                apply_paragraph_style(
                   paragraph, 
                   para_styles[style_idx], 
                   scale=text_scale, 
                   target_language=target_language,
                   font_mapping=font_mapping
                )
            
            if paragraph.runs:
                paragraph.runs[0].font.color.rgb = translated_color

        return True
    except Exception:
        return False


def split_text_chunks(text: str, chunk_size: int) -> list[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return [text]
    chunks = []
    current = []
    current_len = 0
    for line in lines:
        if current_len + len(line) + 1 > chunk_size and current:
            chunks.append("\n".join(current))
            current = []
            current_len = 0
        current.append(line)
        current_len += len(line) + 1
    if current:
        chunks.append("\n".join(current))
    return chunks


def set_corrected_text(
    text_frame: TextFrame, 
    lines: list[str], 
    color: RGBColor | None = None
) -> bool:
    try:
        lines = [sanitize_xml_text(line) for line in lines]
        para_styles = capture_full_frame_styles(text_frame)
        text_frame.clear()
        
        for index, line in enumerate(lines):
            if index == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = line
            
            # Apply original style
            style_idx = min(index, len(para_styles) - 1) if para_styles else -1
            if style_idx >= 0:
                apply_paragraph_style(paragraph, para_styles[style_idx])
            
            # Apply color override
            if color and paragraph.runs:
                for run in paragraph.runs:
                    try:
                        run.font.color.rgb = color
                    except Exception:
                        pass
        return True
    except Exception:
        return False
