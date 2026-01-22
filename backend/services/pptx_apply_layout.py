from __future__ import annotations

import io
from copy import deepcopy
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.slide import Slide
from pptx.text.text import TextFrame


def capture_font_spec(text_frame: TextFrame) -> dict[str, Any]:
    """Captures font specs from the first run of the first paragraph."""
    spec = {
        "name": None,
        "size": None,
        "bold": None,
        "italic": None,
        "underline": None,
        "color": None,
        "alignment": None,
        "level": 0,
        "pPr_xml": None,
    }
    try:
        if text_frame.paragraphs:
            # Find first paragraph with text, or fallback to first
            p = text_frame.paragraphs[0]
            for candidate_p in text_frame.paragraphs:
                if candidate_p.text and candidate_p.text.strip():
                    p = candidate_p
                    break

            spec["alignment"] = p.alignment
            spec["level"] = p.level
            
            # Capture XML for high fidelity
            try:
                from lxml import etree
                p_xml = p._p
                p_pr_xml = p_xml.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pPr')
                if p_pr_xml is None:
                    p_pr_xml = p_xml.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                if p_pr_xml is not None:
                    spec["pPr_xml"] = etree.tostring(p_pr_xml)
            except Exception: pass

            if p.runs:
                run = p.runs[0]
                spec["name"] = run.font.name
                spec["size"] = run.font.size
                spec["bold"] = run.font.bold
                spec["italic"] = run.font.italic
                spec["underline"] = run.font.underline
                try: spec["color"] = run.font.color.rgb
                except Exception: pass
    except Exception:
        pass
    return spec


def add_overflow_textboxes(
    slide: Slide,
    base_shape,
    chunks: list[str],
    font_spec: dict[str, Any] | None,
    translated_color: RGBColor | None = None,
    max_bottom: int = 0,
    target_language: str | None = None,
    font_mapping: dict[str, list[str]] | None = None,
    scale: float = 1.0,
) -> None:
    if not chunks:
        return
    top = base_shape.top
    height = base_shape.height
    margin = int(height * 0.1)
    for index, chunk in enumerate(chunks):
        next_top = top + height + margin + index * (height + margin)
        if next_top + height > max_bottom:
            break
        box = slide.shapes.add_textbox(base_shape.left, next_top, base_shape.width, height)
        text_frame = box.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()
        for line_index, line in enumerate(chunk.split("\n")):
            paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            
            # Apply paragraph-level styles
            # Apply paragraph-level styles
            if font_spec:
                if font_spec.get("alignment") is not None:
                    paragraph.alignment = font_spec["alignment"]
                paragraph.level = font_spec.get("level", 0)
                
                # Apply XML style if available
                if font_spec.get("pPr_xml"):
                    try:
                        from lxml import etree
                        target_p = paragraph._p
                        new_p_pr = etree.fromstring(font_spec["pPr_xml"])
                        
                        # Replace or update existing pPr
                        old_p_pr = target_p.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}pPr')
                        if old_p_pr is None:
                            old_p_pr = target_p.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                        
                        if old_p_pr is not None:
                            target_p.replace(old_p_pr, new_p_pr)
                        else:
                            target_p.insert(0, new_p_pr)
                    except Exception: pass

            if paragraph.runs:
                run = paragraph.runs[0]
                if font_spec:
                    if font_spec.get("name"):
                        run.font.name = font_spec["name"]
                    if font_spec.get("size"):
                        # Apply both global scale and internal scale
                        final_size = font_spec["size"]
                        if scale != 1.0:
                            final_size = int(final_size * scale)
                        run.font.size = final_size
                        
                    if font_spec.get("bold") is not None:
                        run.font.bold = font_spec["bold"]
                    if font_spec.get("italic") is not None:
                        run.font.italic = font_spec["italic"]
                    if font_spec.get("underline") is not None:
                        run.font.underline = font_spec["underline"]
                    
                    # Inherit color from spec if not overridden
                    if font_spec.get("color"):
                        try: run.font.color.rgb = font_spec["color"]
                        except Exception: pass

                # Override with font mapping if needed
                if target_language or font_mapping:
                    from backend.services.font_manager import clone_font_props
                    clone_font_props(run.font, run.font, target_language=target_language, font_mapping=font_mapping)

                if translated_color:
                    try:
                        run.font.color.rgb = translated_color
                    except Exception:
                        pass


def duplicate_slide(presentation: Presentation, slide: Slide) -> tuple[Slide, dict[int, int]]:
    """
    Duplicate a slide with full fidelity including images and shapes.
    
    Uses a hybrid approach:
    - Pictures are re-added using add_picture() with extracted blob data
    - Other shapes are XML-copied with ID regeneration
    
    Returns:
        tuple: (new_slide, shape_id_map) where shape_id_map maps old shape IDs to new shape IDs.
    """
    try:
        import random
        
        print(f"[HYBRID_DUP] Starting duplicate_slide for slide with {len(slide.shapes)} shapes", flush=True)
        
        # 1. Create a new slide based on the same layout
        new_slide = presentation.slides.add_slide(slide.slide_layout)
        
        # Clear any default layout shapes (we'll copy everything from source)
        for shape in list(new_slide.shapes):
            try:
                el = shape.element
                el.getparent().remove(el)
            except Exception as e:
                print(f"[HYBRID_DUP] Failed to clear placeholder: {e}", flush=True)
                pass
        
        shape_id_map = {}
        base_id = random.randint(10000, 99999999)
        id_counter = 0
        
        # Namespaces
        p_ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        
        # 2. Copy relationships (excluding layout and master)
        rid_mapping = {}
        for rel in slide.part.rels.values():
            if rel.reltype in (
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/master"
            ):
                continue
            
            old_rid = rel.rId
            try:
                if rel.is_external:
                    new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
                
                rid_mapping[old_rid] = new_rid
            except Exception as e:
                print(f"[XML_DUP] Failed to copy relationship {old_rid} ({rel.reltype}): {e}", flush=True)
        
        # 3. Process each shape
        sp_tree = new_slide.shapes._spTree
        ext_lst = sp_tree.find(f"{p_ns}extLst")
        
        for shape in slide.shapes:
            try:
                old_id = shape.shape_id
                shape_name = getattr(shape, "name", "Unknown")
                
                # Use unified XML deep copy for ALL shapes
                new_element = deepcopy(shape.element)
                
                # Pass 1: Regenerate all shape IDs
                for elem in new_element.iter():
                    if elem.tag.endswith('}cNvPr'):
                        try:
                            current_id = elem.get('id')
                            if current_id:
                                old_sid = int(current_id)
                                new_id = base_id + id_counter
                                id_counter += 1
                                elem.set('id', str(new_id))
                                shape_id_map[old_sid] = new_id
                        except ValueError:
                            pass
                
                # Pass 2: Fix references (Connectors and Relationships)
                for elem in new_element.iter():
                    # Fix connector references
                    if elem.tag in (f"{a_ns}stCxn", f"{a_ns}endCxn"):
                        ref_id = elem.get('id')
                        if ref_id and ref_id.isdigit():
                            rid_int = int(ref_id)
                            if rid_int in shape_id_map:
                                elem.set('id', str(shape_id_map[rid_int]))
                    
                    # Fix relationship references (r:id, r:embed, r:link, etc.)
                    for attr_name, attr_val in list(elem.attrib.items()):
                        if 'http://schemas.openxmlformats.org/officeDocument/2006/relationships' in attr_name:
                            if attr_val in rid_mapping:
                                elem.attrib[attr_name] = rid_mapping[attr_val]
                
                # Insert into shape tree
                if ext_lst is not None:
                    ext_lst.addprevious(new_element)
                else:
                    sp_tree.append(new_element)
                    
            except Exception as e:
                print(f"[HYBRID_DUP] Failed to copy shape ID {shape.shape_id}: {e}", flush=True)
                continue
        
        return new_slide, shape_id_map

    except Exception as exc:
        print(f"[HYBRID_DUP] CRITICAL: duplicate_slide completely failed. Falling back to Blank slide. Error: {exc}", flush=True)
        return presentation.slides.add_slide(presentation.slide_layouts[6]), {}


def insert_slide_after(presentation: Presentation, new_slide: Slide, after_index: int) -> None:
    sld_id_list = presentation.slides._sldIdLst
    new_sld_id = sld_id_list[-1]
    sld_id_list.remove(new_sld_id)
    sld_id_list.insert(after_index + 1, new_sld_id)


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
        except Exception:
            continue


def find_shape_in_shapes(shapes, shape_id: int):
    for shape in _iter_shapes(shapes):
        if shape.shape_id == shape_id:
            return shape
    return None


def find_shape_with_id(slide: Slide, shape_id: int):
    return find_shape_in_shapes(slide.shapes, shape_id)


def iter_table_cells(shape) -> list[TextFrame]:
    cells = []
    for row in shape.table.rows:
        for cell in row.cells:
            cells.append(cell.text_frame)
    return cells

def fix_title_overlap(slide: Slide) -> None:
    """
    Detects if the title overlaps with other content.
    Groups vertically aligned obstacles (e.g. stack of images) and moves/scales them together
    to fit between the Title and any Bottom Limit (e.g. Table/Footer).
    """
    try:
        title_shape = None
        if slide.shapes.title:
            title_shape = slide.shapes.title
        else:
            # Fallback: top-most text box
            sorted_shapes = sorted(
                [s for s in slide.shapes if s.has_text_frame], 
                key=lambda s: s.top
            )
            if sorted_shapes:
                title_shape = sorted_shapes[0]
        
        if not title_shape: return

        # Title boundary with buffer
        margin_buffer = 50000 
        title_bottom = title_shape.top + title_shape.height + margin_buffer
        slide_height = 6858000 
        try: slide_height = slide.part.presentation.slide_height
        except: pass

        # 1. Identify valid obstacles
        valid_types = (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP, 
                       MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.AUTO_SHAPE)
        
        obstacles = [
            s for s in slide.shapes 
            if s.shape_id != title_shape.shape_id 
            and s.shape_type in valid_types
        ]

        if not obstacles: return

        # 2. Group obstacles into "Columns" based on X-overlap
        # Simple clustering: if horizontal overlap > 0.5 * min_width
        columns = []
        sorted_obs = sorted(obstacles, key=lambda s: s.left)
        
        while sorted_obs:
            current = sorted_obs.pop(0)
            col = [current]
            
            # Find all subsequent shapes that align with 'current'
            # We iterate a copy to safely modify sorted_obs
            remaining = []
            current_center = current.left + current.width / 2
            
            for other in sorted_obs:
                other_center = other.left + other.width / 2
                # Check alignment: centers are close (within 20% of width)
                # or significant overlap
                overlap_x = min(current.left + current.width, other.left + other.width) - max(current.left, other.left)
                is_aligned = False
                
                if overlap_x > 0:
                   share_pct = overlap_x / min(current.width, other.width)
                   if share_pct > 0.5: is_aligned = True
                
                if is_aligned:
                    col.append(other)
                else:
                    remaining.append(other)
            
            columns.append(col)
            sorted_obs = remaining

        # 3. Process each column
        for col in columns:
            # Sort column by top
            col.sort(key=lambda s: s.top)
            
            top_shape = col[0]
            bottom_shape = col[-1]
            col_top = top_shape.top
            col_bottom = bottom_shape.top + bottom_shape.height
            
            # Check overlap with Title
            if col_top < title_bottom and col_bottom > title_shape.top:
                # We have a collision.
                
                # Detect Floor (Limit Bottom) for this specific column
                limit_bottom = slide_height
                
                # Check against ALL other shapes not in this column
                col_ids = {s.shape_id for s in col}
                for potential in slide.shapes:
                    if potential.shape_id in col_ids: continue
                    if potential.shape_id == title_shape.shape_id: continue
                    
                    # Must be below current column bottom? No, below current column TOP (potential block)
                    # Actually we want the nearest object strictly below the "would-be" position?
                    # Let's find obstacles strictly below the current column footprint
                    if potential.top >= col_top: # Candidate for floor
                         # Check horizontal overlap with column
                         c_left = min(s.left for s in col)
                         c_width = max(s.left + s.width for s in col) - c_left
                         
                         p_overlap = min(c_left + c_width, potential.left + potential.width) - max(c_left, potential.left)
                         if p_overlap > 0:
                             if potential.top < limit_bottom:
                                 limit_bottom = potential.top

                # Calculate Available Space
                margin = 50000
                target_top = title_bottom + margin
                max_frame_height = (limit_bottom - margin) - target_top
                
                if max_frame_height < 100000:
                    # Too tight, just nudge top
                     pass 
                else:
                    # Current Height of the group
                    current_group_height = col_bottom - col_top
                    
                    # Determine Scale Factor
                    scale = 1.0
                    if current_group_height > max_frame_height:
                        scale = max_frame_height / current_group_height
                    
                    # Apply Transformation
                    # 1. Calculate new top for the first element
                    # 2. Re-space subsequent elements relative to first
                    
                    base_top = target_top
                    
                    # We need to maintain relative offsets
                    # offsets[i] = (shape.top - col_top) * scale
                    
                    for s in col:
                        rel_y = s.top - col_top
                        new_rel_y = int(rel_y * scale)
                        
                        s.top = base_top + new_rel_y
                        
                        # Scale size
                        s.height = int(s.height * scale)
                        s.width = int(s.width * scale)
                        
                        # Fix X center? 
                        # Usually scaling is centered or top-left. 
                        # To be safe, keep left as is (or scale relative to center?)
                        # User asked for "Scale Down", usually implies shrinking in place or towards center.
                        # But simple W/H scale + Top adjust is usually safest for layout preservation.
                        pass

    except Exception as e:
        print(f"[LAYOUT_FIX] Error in fix_title_overlap: {e}", flush=True)
