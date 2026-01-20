from __future__ import annotations

from collections.abc import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.text.text import TextFrame

from backend.services.pptx_apply_layout import (
    add_overflow_textboxes,
    capture_font_spec,
    duplicate_slide,
    find_shape_in_shapes,
    find_shape_with_id,
    insert_slide_after,
    iter_table_cells,
)

from backend.services.font_manager import estimate_scale
from backend.services.pptx_apply_text import (
    apply_shape_highlight,
    build_corrected_lines,
    parse_dash_style,
    parse_hex_color,
    set_bilingual_text,
    set_corrected_text,
    set_text_preserve_format,
    split_text_chunks,
)
from backend.services.pptx_xml_core import get_pptx_theme_summary


def _apply_translations_to_presentation(
    presentation: Presentation,
    blocks: Iterable[dict],
    mode: str = "direct",
    target_language: str | None = None,
    font_mapping: dict[str, list[str]] | None = None,
) -> None:
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}
    supported_types = {"textbox", "table_cell", "notes"}

    # Analyze theme once per presentation
    pptx_path = getattr(presentation, "_pptx_path", None) # May need better way to get path
    theme_data = get_pptx_theme_summary(pptx_path) if pptx_path else None

    # Default accent color
    accent_color = RGBColor(0x1F, 0x77, 0xB4)
    if theme_data and 'colors' in theme_data:
        accent1 = theme_data['colors'].get('accent1')
        if accent1:
            try:
                accent_color = RGBColor.from_string(accent1)
            except Exception: pass

    for block in blocks:
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        translated_text = block.get("translated_text", "")
        source_text = block.get("source_text", "")

        if mode == "bilingual":
            combined_text = f"{source_text}\n\n{translated_text}"
        else:
            combined_text = translated_text

        block_type = block.get("block_type", "textbox")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]
        scale = estimate_scale(source_text, translated_text)

        is_auto_layout = (block.get("layout", "auto") in ("auto", "new_slide"))

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = find_shape_in_shapes(slide.notes_slide.shapes, shape_id)
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            if mode == "bilingual":
                if not set_bilingual_text(
                    notes_shape.text_frame,
                    source_text,
                    translated_text,
                    auto_size=is_auto_layout,
                    scale=scale,
                    theme_data=theme_data,
                    target_language=target_language,
                    font_mapping=font_mapping,
                ):
                    set_text_preserve_format(notes_shape.text_frame, combined_text, auto_size=is_auto_layout, scale=scale)
            else:
                set_text_preserve_format(notes_shape.text_frame, translated_text, auto_size=is_auto_layout, scale=scale)
            continue

        # 3. Apply Text
        shape = find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue

            if mode == "bilingual":
                if not set_bilingual_text(
                    cells[idx],
                    source_text,
                    translated_text,
                    auto_size=is_auto_layout,
                    scale=scale,
                    theme_data=theme_data,
                    target_language=target_language,
                    font_mapping=font_mapping,
                ):
                    set_text_preserve_format(cells[idx], combined_text, auto_size=is_auto_layout, scale=scale)
            else:
                set_text_preserve_format(cells[idx], translated_text, auto_size=is_auto_layout, scale=scale)
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue

        # Backup original geometry to prevent collapse
        orig_left = shape.left
        orig_top = shape.top
        orig_width = shape.width
        orig_height = shape.height

        # Handle Overflow
        # In 'direct' mode (New Page), we generally want to avoid splitting if the user wants to stay on one page.
        # But for extremely long text, we still split.
        overflow_limit = 1000 if block.get("layout", "auto") == "new_slide" else 400

        if is_auto_layout and len(translated_text) > overflow_limit:
            if mode == "bilingual":
                set_text_preserve_format(shape.text_frame, source_text, auto_size=is_auto_layout)
            else:
                shape.text_frame.clear()

            chunks = split_text_chunks(translated_text, 300)
            font_spec = capture_font_spec(shape.text_frame)
            max_bottom = presentation.slide_height
            add_overflow_textboxes(
                slide,
                shape,
                chunks,
                font_spec,
                accent_color,
                max_bottom,
            )
            # Restore geometry even if cleared
            shape.left, shape.top, shape.width, shape.height = orig_left, orig_top, orig_width, orig_height
            continue

        if mode == "bilingual":
            if not set_bilingual_text(
                shape.text_frame,
                source_text,
                translated_text,
                auto_size=is_auto_layout,
                scale=scale,
                theme_data=theme_data,
                target_language=target_language,
                font_mapping=font_mapping,
            ):
                set_text_preserve_format(shape.text_frame, combined_text, auto_size=is_auto_layout, scale=scale)
        else:
            if is_auto_layout:
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            set_text_preserve_format(shape.text_frame, translated_text, auto_size=is_auto_layout, scale=scale)

        # Restore original geometry to prevent any unintended shrinking/movement
        try:
            shape.left = orig_left
            shape.top = orig_top
            shape.width = orig_width
            shape.height = orig_height
        except Exception:
            pass


def apply_translations(
    pptx_in: str,
    pptx_out: str,
    blocks: Iterable[dict],
    mode: str = "direct",
    target_language: str | None = None,
    font_mapping: dict[str, list[str]] | None = None,
) -> None:
    presentation = Presentation(pptx_in)
    presentation._pptx_path = pptx_in # For theme analysis
    _apply_translations_to_presentation(
        presentation,
        blocks,
        mode=mode,
        target_language=target_language,
        font_mapping=font_mapping,
    )
    presentation.save(pptx_out)


def apply_bilingual(
    pptx_in: str,
    pptx_out: str,
    blocks: list[dict],
    layout: str = "inline",
    target_language: str | None = None,
    font_mapping: dict[str, list[str]] | None = None,
) -> None:
    presentation = Presentation(pptx_in)
    presentation._pptx_path = pptx_in # For theme analysis
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}
    supported_types = {"textbox", "table_cell", "notes"}

    if layout == "new_slide":
        print(f"[APPLY_BILINGUAL] Using Hybrid Duplication for new_slide mode", flush=True)
        slide_blocks: dict[int, list[dict]] = {}
        for block in blocks:
            if block.get("apply") is False or block.get("selected") is False:
                continue
            slide_index = block.get("slide_index")
            if slide_index is None:
                continue
            slide_blocks.setdefault(slide_index, []).append(block)

        new_blocks = []
        offset = 0
        for slide_index in range(len(presentation.slides)):
            if slide_index not in slide_blocks:
                continue
            source_slide = presentation.slides[slide_index + offset]
            try:
                res = duplicate_slide(presentation, source_slide)
                if isinstance(res, tuple):
                    new_slide, shape_map = res
                else:
                    new_slide = res
                    shape_map = {}
            except Exception as e:
                print(f"[APPLY_BILINGUAL] duplicate_slide failed: {e}", flush=True)
                new_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                shape_map = {}

            insert_slide_after(presentation, new_slide, slide_index + offset)
            offset += 1
            new_slide_index = slide_index + offset

            print(f"[APPLY_BILINGUAL] Duplicated slide {slide_index} -> new slide at {new_slide_index}", flush=True)

            for block in slide_blocks[slide_index]:
                updated = dict(block)
                updated["slide_index"] = new_slide_index

                # Remap shape_id if possible
                old_sid = updated.get("shape_id")
                if old_sid in shape_map:
                    updated["shape_id"] = shape_map[old_sid]
                    print(f"[APPLY_BILINGUAL] Remapped shape_id {old_sid} -> {shape_map[old_sid]}", flush=True)

                new_blocks.append(updated)

        print(f"[APPLY_BILINGUAL] Applying translations to {len(new_blocks)} shapes in new slides with layout={layout}", flush=True)
        presentation._pptx_path = pptx_in # Re-inject path
        _apply_translations_to_presentation(
            presentation,
            new_blocks,
            mode="direct",
            target_language=target_language,
            font_mapping=font_mapping,
        )
        presentation.save(pptx_out)
        return

    # Prepare theme color for translated text
    theme_data = get_pptx_theme_summary(pptx_in)
    translated_color = RGBColor(0x1F, 0x77, 0xB4) # Default Blue
    if theme_data and "colors" in theme_data and "accent1" in theme_data["colors"]:
        try:
            translated_color = parse_hex_color(theme_data["colors"]["accent1"])
        except Exception:
            pass

    for block in blocks:
        if block.get("apply") is False or block.get("selected") is False:
            continue
        block_type = block.get("block_type")
        if block_type not in supported_types:
            continue
        translated_text = block.get("translated_text", "")
        if not translated_text:
            continue
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]
        source_text = block.get("source_text", "")
        # Extract target language from block metadata if available, though typically blocks don't store "target_language" directly unless enriched.
        # But wait, the standard block structure doesn't include target_language.
        # However, looking at previous conversations, "fixing translation language code", maybe it's passed?
        # Actually, "blocks" usually just have text. The 'target_language' is usually global to the request.
        # BUT, if we want to support this, we need to know the target language.
        # In apply_translate, target_language is not readily available in 'blocks'.
        # Assuming we might need to rely on detection or assume 'vi' based on context?
        # No, the USER REQUEST context implies we know it.
        # Checking backend/api/pptx.py, apply_bilingual receives blocks_data.
        # It does NOT receive target_language as an argument.
        # This is a limitation. I should update apply_bilingual signature or use a heuristic.
        # Wait, the user prompt showed "vi" in the screenshot.
        # Let's inspect block structure more deeply. Maybe it's not there.
        # If not present, I can't pass it.
        # But wait, I can assume the user wants this feature.
        # Let's check where apply_bilingual is called from. api/pptx.py.
        # It's called from `pptx_apply`. `pptx_apply` payload does NOT have target_language.

        # PLAN ADJUSTMENT: I need to update `pptx_apply` API to accept `target_language` form field.
        # Then pass it to apply_bilingual.

        # For now, let's inject it into set_bilingual_text call, expecting it MIGHT be in block for now?
        # Or better, let's update call sites.

        # Let's look at `block` structure again.
        # `translated_text` is there.
        # I will extract `target_language` from block if it exists (maybe added by translate step?)
        # If not, I will add it to the API.

        target_language = block.get("target_language")

        combined_text = f"{source_text}\\n\\n{translated_text}"
        scale = estimate_scale(source_text, translated_text)

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = find_shape_in_shapes(slide.notes_slide.shapes, shape_id)
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            if not set_bilingual_text(
                notes_shape.text_frame,
                source_text,
                translated_text,
                auto_size=layout == "auto",
                scale=scale,
                theme_data=theme_data,
                target_language=target_language,
                font_mapping=font_mapping,
            ):
                set_text_preserve_format(notes_shape.text_frame, combined_text, scale=scale)
            continue

        shape = find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            if not set_bilingual_text(
                cells[idx],
                source_text,
                translated_text,
                auto_size=layout == "auto",
                scale=scale,
                theme_data=theme_data,
                target_language=target_language,
                font_mapping=font_mapping,
            ):
                set_text_preserve_format(cells[idx], combined_text, scale=scale)
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        if layout == "auto" and len(translated_text) > 400:
            set_text_preserve_format(shape.text_frame, source_text)
            chunks = split_text_chunks(translated_text, 300)
            font_spec = capture_font_spec(shape.text_frame)

            max_bottom = presentation.slide_height
            add_overflow_textboxes(
                slide,
                shape,
                chunks,
                font_spec,
                translated_color,
                max_bottom,
            )
            continue
        if not set_bilingual_text(
            shape.text_frame,
            source_text,
            translated_text,
            auto_size=layout == "auto",
            scale=scale,
            theme_data=theme_data,
            target_language=target_language,
            font_mapping=font_mapping,
        ):
            set_text_preserve_format(shape.text_frame, combined_text, scale=scale)

    presentation.save(pptx_out)


def apply_chinese_corrections(
    pptx_in: str,
    pptx_out: str,
    blocks: list[dict],
    fill_color: str | None = None,
    text_color: str | None = None,
    line_color: str | None = None,
    line_dash: str | None = None,
) -> None:
    presentation = Presentation(pptx_in)
    table_cell_positions: dict[tuple[int, int], list[TextFrame]] = {}
    table_cell_index: dict[tuple[int, int], int] = {}
    supported_types = {"textbox", "table_cell", "notes"}
    fill_yellow = parse_hex_color(fill_color, RGBColor(0xFF, 0xF1, 0x6A))
    text_red = parse_hex_color(text_color, RGBColor(0xD9, 0x00, 0x00))
    line_purple = parse_hex_color(line_color, RGBColor(0x7B, 0x2C, 0xB9))
    dash_style = parse_dash_style(line_dash) or MSO_LINE_DASH_STYLE.DASH

    for block in blocks:
        if block.get("apply") is False or block.get("selected") is False:
            continue
        block_type = block.get("block_type")
        if block_type not in supported_types:
            continue
        translated_text = block.get("translated_text", "")
        if not translated_text:
            continue
        slide_index = block.get("slide_index")
        shape_id = block.get("shape_id")
        if slide_index is None or shape_id is None:
            continue
        if slide_index < 0 or slide_index >= len(presentation.slides):
            continue
        slide = presentation.slides[slide_index]
        source_text = block.get("source_text", "")
        lines = build_corrected_lines(source_text, translated_text)

        if block_type == "notes":
            if not slide.has_notes_slide:
                continue
            notes_shape = find_shape_in_shapes(slide.notes_slide.shapes, shape_id)
            if notes_shape is None or not notes_shape.has_text_frame:
                continue
            apply_shape_highlight(notes_shape, fill_yellow, line_purple, dash_style)
            if not set_corrected_text(notes_shape.text_frame, lines, text_red):
                set_text_preserve_format(notes_shape.text_frame, "\n".join(lines))
            continue

        shape = find_shape_with_id(slide, shape_id)
        if shape is None:
            continue

        if block_type == "table_cell":
            if not shape.has_table:
                continue
            key = (slide_index, shape_id)
            if key not in table_cell_positions:
                table_cell_positions[key] = iter_table_cells(shape)
                table_cell_index[key] = 0
            idx = table_cell_index[key]
            cells = table_cell_positions[key]
            if idx >= len(cells):
                continue
            try:
                table = shape.table
                col_count = len(table.columns)
                row_index = idx // col_count
                col_index = idx % col_count
                cell = table.cell(row_index, col_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_yellow
            except Exception:
                pass
            apply_shape_highlight(shape, fill_yellow, line_purple, dash_style)
            if not set_corrected_text(cells[idx], lines, text_red):
                set_text_preserve_format(cells[idx], "\n".join(lines))
            table_cell_index[key] = idx + 1
            continue

        if not shape.has_text_frame:
            continue
        apply_shape_highlight(shape, fill_yellow, line_purple, dash_style)
        if not set_corrected_text(shape.text_frame, lines, text_red):
            set_text_preserve_format(shape.text_frame, "\n".join(lines))

    presentation.save(pptx_out)
