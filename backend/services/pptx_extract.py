from __future__ import annotations

import json
import re
from collections.abc import Iterable
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide
from pptx.table import _Cell
from pptx.text.text import TextFrame

from backend.contracts import make_block


# Load preserve terms at module level for performance
def _load_preserve_terms() -> list[dict]:
    """Load preserve terms from JSON file."""
    preserve_file = Path(__file__).parent.parent / "data" / "preserve_terms.json"
    if not preserve_file.exists():
        return []
    try:
        with open(preserve_file, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []


PRESERVE_TERMS = _load_preserve_terms()


def _is_numeric_only(text: str) -> bool:
    """Check if the text consists only of numbers, punctuation, or whitespace."""
    if not text.strip():
        return True
    # If it contains any letter (English, CJK), it's not purely numeric/symbolic
    if re.search(r"[a-zA-Z\u4e00-\u9fff\u3040-\u30ff\uac00-\ud7af]", text):
        return False
    return True


def _is_technical_terms_only(text: str) -> bool:
    """
    Check if the text consists only of technical terms, product names, or acronyms.
    First checks against preserve_terms.json database, then falls back to auto-detection.
    """
    if not text.strip():
        return True

    # Priority 1: Check preserve terms database
    text_clean = text.strip()
    for term_entry in PRESERVE_TERMS:
        term = term_entry.get("term", "")
        case_sensitive = term_entry.get("case_sensitive", True)

        if case_sensitive:
            if text_clean == term:
                return True
        else:
            if text_clean.lower() == term.lower():
                return True

    # Priority 2: Auto-detection fallback
    # Remove common separators
    cleaned = re.sub(r"[,、，/\s]+", " ", text).strip()

    # Check if contains any CJK characters (if yes, it's not pure technical terms)
    if re.search(r"[\u4e00-\u9fff\u3040-\u30ff\u0e00-\u0e7f]", cleaned):
        return False

    # Check if contains sentence-forming words (articles, prepositions, verbs)
    sentence_indicators = (
        r"\b(the|a|an|is|are|was|were|be|have|has|had|do|does|did|will|would|can|could|should|"
        r"may|might|must|please|this|that|these|those|with|from|to|in|on|at|for|of|and|or|but)\b"
    )
    if re.search(sentence_indicators, cleaned, re.IGNORECASE):
        return False

    # Split into words
    words = cleaned.split()
    word_count = len(words)

    # Heuristic: True technical term lists are usually short (1-3 words) 
    # and match specific patterns (ALLCAPS, acronyms, specific product names).
    # Full sentences with Title Case (e.g., "Welcome to Our System") should NOT be filtered.
    if word_count > 10:
        return False
        
    # Patterns: 
    # 1. ALL UPPERCASE (SOP, CRM, API)
    # 2. MixedCase (GraphQL, iOS, macOS)
    # 3. Simple TitleCase (Only if 1-3 words, likely a Label/Product)
    is_all_caps = all(re.match(r"^[A-Z0-9_\-]+$", w) for w in words)
    is_mixed_case = all(re.match(r"^[A-Z][a-z]*[A-Z][a-zA-Z]*$", w) for w in words)
    is_title_case = all(re.match(r"^[A-Z][a-z]+$", w) for w in words)
    is_pure_lower = all(re.match(r"^[a-z]+$", w) for w in words)

    if is_all_caps or is_mixed_case:
        return True
    
    # TitleCase or pure lower is only filtered if very short (likely a UI label or single term)
    if (is_title_case or is_pure_lower) and word_count <= 1:
        return True

    return False


def _text_frame_to_text(text_frame: TextFrame) -> str:
    paragraphs = [paragraph.text for paragraph in text_frame.paragraphs]
    return "\n".join(paragraphs).strip()


def _iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
        except Exception:
            continue


def emu_to_points(emu: int | float) -> float:
    """Convert EMU to Points (72 points per inch). 1 Point = 12700 EMUs."""
    return float(emu) / 12700.0


def _iter_textbox_blocks(slide: Slide, slide_index: int, seen_ids: set[int]) -> Iterable[dict]:
    for shape in _iter_shapes(slide.shapes):
        try:
            if not shape.has_text_frame or shape.has_table:
                continue
            if shape.shape_id in seen_ids:
                continue
            text = _text_frame_to_text(shape.text_frame)
        except Exception:
            continue
        if not text or _is_numeric_only(text) or _is_technical_terms_only(text):
            continue
        
        seen_ids.add(shape.shape_id)
        # Extract layout info
        x = emu_to_points(shape.left)
        y = emu_to_points(shape.top)
        w = emu_to_points(shape.width)
        h = emu_to_points(shape.height)

        yield make_block(slide_index, shape.shape_id, "textbox", text, x=x, y=y, width=w, height=h)


def _cell_to_text(cell: _Cell) -> str:
    return _text_frame_to_text(cell.text_frame)


def _iter_table_blocks(slide: Slide, slide_index: int) -> Iterable[dict]:
    for shape in _iter_shapes(slide.shapes):
        if not shape.has_table:
            continue
        
        # Table position (same for all cells in this table for simplification in preview)
        tx = emu_to_points(shape.left)
        ty = emu_to_points(shape.top)
        tw = emu_to_points(shape.width)
        th = emu_to_points(shape.height)

        for row in shape.table.rows:
            for cell in row.cells:
                text = _cell_to_text(cell)
                if not text or _is_numeric_only(text) or _is_technical_terms_only(text):
                    continue
                yield make_block(slide_index, shape.shape_id, "table_cell", text, x=tx, y=ty, width=tw, height=th)


def _iter_notes_blocks(slide: Slide, slide_index: int) -> Iterable[dict]:
    if not slide.has_notes_slide:
        return
    for shape in _iter_shapes(slide.notes_slide.shapes):
        try:
            if not shape.has_text_frame:
                continue
            text = _text_frame_to_text(shape.text_frame)
        except Exception:
            continue
        if not text or _is_numeric_only(text) or _is_technical_terms_only(text):
            continue
        
        # Notes position
        x = emu_to_points(shape.left)
        y = emu_to_points(shape.top)
        w = emu_to_points(shape.width)
        h = emu_to_points(shape.height)

        yield make_block(slide_index, shape.shape_id, "notes", text, x=x, y=y, width=w, height=h)


def extract_blocks(pptx_path: str) -> dict:
    presentation = Presentation(pptx_path)
    blocks: list[dict] = []
    
    # Get slide dimensions in Points
    slide_width = emu_to_points(presentation.slide_width)
    slide_height = emu_to_points(presentation.slide_height)

    for slide_index, slide in enumerate(presentation.slides):
        seen_ids = set()
        blocks.extend(_iter_textbox_blocks(slide, slide_index, seen_ids))
        blocks.extend(_iter_table_blocks(slide, slide_index))
        blocks.extend(_iter_notes_blocks(slide, slide_index))
    
    return {
        "blocks": blocks,
        "slide_width": slide_width,
        "slide_height": slide_height
    }

