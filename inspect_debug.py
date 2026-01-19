from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_pptx(path, slide_idx):
    pp = Presentation(path)
    slide = pp.slides[slide_idx]
    print(f"Slide index {slide_idx} (Slide {slide_idx+1}) total shapes: {len(slide.shapes)}")
    
    for shape in slide.shapes:
        info = {
            "id": shape.shape_id,
            "type": str(shape.shape_type),
            "name": shape.name,
            "has_text": shape.has_text_frame,
            "text": shape.text_frame.text[:100] if shape.has_text_frame else "",
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height
        }
        
        fill_info = "None"
        try:
            if hasattr(shape, "fill") and shape.fill.type == 1:
                fill_info = str(shape.fill.fore_color.rgb)
            elif hasattr(shape, "fill"):
                fill_info = str(shape.fill.type)
        except:
            fill_info = "Error"
            
        print(f"ID:{info['id']} | Type:{info['type']} | Name:{info['name']} | Fill:{fill_info} | Text:[{info['text'].replace('\n', ' ')}]")

import sys
import io

if __name__ == "__main__":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    path = r'D:\Downloads\pptx_translated (7).pptx'
    print("=== SLIDE 27 (Original) ===")
    inspect_pptx(path, 26)
    print("\n=== SLIDE 28 (Duplicated) ===")
    inspect_pptx(path, 27)
